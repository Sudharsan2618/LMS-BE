from datetime import datetime
from typing import List, Dict, Optional
import os
import tempfile
import boto3
from botocore.exceptions import ClientError
from dotenv import load_dotenv

from psycopg2.extras import RealDictCursor
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from app.utils.db_utils import get_db_connection
from app.config.database import DB_CONFIG

# Load environment variables
load_dotenv()


def _truncate_text(text: str, max_chars: int) -> str:
    if text is None:
        return ""
    text = text.strip()
    return text if len(text) <= max_chars else text[: max_chars - 3].rstrip() + "..."


def _add_bullets_to_placeholder(placeholder, bullets: List[str], font_size_pt: int = 20) -> None:
    text_frame = placeholder.text_frame
    text_frame.clear()
    for idx, bullet in enumerate(bullets):
        if not bullet.strip():  # Skip empty bullets
            continue
        p = text_frame.add_paragraph() if idx > 0 else text_frame.paragraphs[0]
        p.text = bullet
        p.level = 0
        run = p.runs[0]
        run.font.size = Pt(font_size_pt)
        p.alignment = PP_ALIGN.LEFT


def _get_layout(prs: Presentation, name_fallback_index: int = 1):
    # Try to pick a Title and Content layout; fallback to index if names vary
    for layout in prs.slide_layouts:
        if layout.name and ("Title and Content" in layout.name or "Title and Body" in layout.name):
            return layout
    return prs.slide_layouts[name_fallback_index]


def update_course_ppt_url(conn, course_id: int, ppt_url: str) -> bool:
    """
    Update the PPT URL in the course_master table
    
    :param conn: Database connection
    :param course_id: ID of the course to update
    :param ppt_url: S3 URL of the generated PPT
    :return: True if successful, False otherwise
    """
    try:
        with conn.cursor() as cursor:
            cursor.execute(
                "UPDATE lms.course_master SET ppt_url = %s WHERE course_id = %s",
                (ppt_url, course_id)
            )
            
            # Check if any row was updated
            if cursor.rowcount > 0:
                conn.commit()
                print(f"✅ Database updated: PPT URL saved for course_id {course_id}")
                return True
            else:
                print(f"⚠️  No course found with course_id {course_id}")
                return False
                
    except Exception as e:
        print(f"❌ Failed to update database: {str(e)}")
        conn.rollback()
        return False


def upload_to_s3(file_path: str, bucket_name: str, object_name: str) -> str:
    """
    Upload a file to an S3 bucket and return the public URL
    
    :param file_path: File to upload
    :param bucket_name: Bucket to upload to
    :param object_name: S3 object name
    :return: Public URL of uploaded file or None if error
    """
    # Get AWS credentials from environment
    aws_access_key = os.getenv('AWS_ACCESS_KEY_ID')
    aws_secret_key = os.getenv('AWS_SECRET_ACCESS_KEY')
    aws_region = os.getenv('AWS_REGION', 'us-east-1')
    
    if not aws_access_key or not aws_secret_key:
        raise ValueError("AWS credentials not found in environment variables")
    
    # Create S3 client
    s3_client = boto3.client(
        's3',
        aws_access_key_id=aws_access_key,
        aws_secret_access_key=aws_secret_key,
        region_name=aws_region
    )
    
    try:
        # Upload the file
        s3_client.upload_file(
            file_path, 
            bucket_name, 
            object_name,
            ExtraArgs={
                'ContentType': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'ACL': 'public-read'  # Make file publicly accessible
            }
        )
        
        # Return the public URL
        public_url = f"https://{bucket_name}.s3.{aws_region}.amazonaws.com/{object_name}"
        return public_url
        
    except ClientError as e:
        print(f"Error uploading to S3: {e}")
        raise


def generate_ppt_for_course(course_id: int, filename: Optional[str] = None, template_path: Optional[str] = None, max_slides: int = 30, upload_to_cloud: bool = True) -> Dict[str, str]:
    """
    Generate a PPTX for a course using stored content from DB and optionally upload to S3.

    :param course_id: ID of the course to generate PPT for
    :param filename: Optional custom filename (without extension)
    :param template_path: Optional path to PowerPoint template
    :param max_slides: Maximum number of slides to generate
    :param upload_to_cloud: Whether to upload to S3 (default: True)
    :return: Dictionary with 'local_path', 'cloud_url', and 'filename'
    """
    conn = None
    try:
        conn = get_db_connection(DB_CONFIG)
        if not conn:
            raise RuntimeError("Database connection failed")

        query = (
            """
            SELECT course_id, course_mastertitle_breakdown_id, course_mastertitle_breakdown,
                   course_subtitle_id, course_subtitle, subtitle_content, subtitle_code, subtitle_help_text,
                   helpfull_links
            FROM lms.course_content
            WHERE course_id = %s
            ORDER BY course_mastertitle_breakdown_id, course_subtitle_id
            """
        )
        with conn.cursor(cursor_factory=RealDictCursor) as cursor:
            cursor.execute(query, (course_id,))
            rows: List[Dict] = cursor.fetchall()

        if not rows:
            raise ValueError(f"No content found for course_id={course_id}")

        # Build comprehensive course structure
        course_title = rows[0].get("course_mastertitle_breakdown") or f"Course {course_id}"
        sections: Dict[int, Dict] = {}
        for r in rows:
            sect_id = r["course_mastertitle_breakdown_id"]
            sections.setdefault(sect_id, {"title": r["course_mastertitle_breakdown"], "items": []})
            sections[sect_id]["items"].append(r)
        
        print(f"Found {len(sections)} sections with {len(rows)} total content items")

        prs = Presentation(template_path) if template_path else Presentation()
        title_layout = prs.slide_layouts[0]
        content_layout = _get_layout(prs)

        # Slide 1: Title
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = _truncate_text(course_title, 100)
        if len(slide.placeholders) > 1:
            subtitle_ph = slide.placeholders[1]
            subtitle_ph.text = f"Comprehensive Learning Guide\nCourse ID: {course_id} | {len(sections)} Modules | {len(rows)} Topics"

        # Add course overview slide
        if len(sections) > 1:
            overview_slide = prs.slides.add_slide(content_layout)
            overview_slide.shapes.title.text = "Course Overview"
            if len(overview_slide.placeholders) > 1:
                body = overview_slide.placeholders[1]
                overview_bullets = []
                for sect_id in sorted(sections.keys()):
                    section = sections[sect_id]
                    section_title = section["title"] or f"Module {sect_id}"
                    item_count = len(section["items"])
                    overview_bullets.append(f"{section_title} ({item_count} topics)")
                _add_bullets_to_placeholder(body, overview_bullets, font_size_pt=18)
        
        # Slides 2..N: content
        slide_count = 2 if len(sections) > 1 else 1
        for sect_id in sorted(sections.keys()):
            section = sections[sect_id]
            section_title = section["title"] or f"Module {sect_id}"
            
            # Add section header slide
            if slide_count < max_slides - 2:
                s = prs.slides.add_slide(content_layout)
                slide_count += 1
                s.shapes.title.text = f"📚 {_truncate_text(section_title, 70)}"
                if len(s.placeholders) > 1:
                    body = s.placeholders[1]
                    bullets = []
                    bullets.append(f"📊 Topics covered: {len(section['items'])}")
                    bullets.append("🎯 Learning objectives:")
                    
                    # Add first few subtitles as objectives
                    for item in section["items"][:5]:
                        bullet_text = item.get("course_subtitle") or ""
                        if bullet_text:
                            bullets.append(f"  • {_truncate_text(bullet_text, 80)}")
                    _add_bullets_to_placeholder(body, bullets, font_size_pt=18)

            # Add detailed content slides for each item
            for item in section["items"]:
                if slide_count >= max_slides - 2:
                    break
                s = prs.slides.add_slide(content_layout)
                slide_count += 1
                sub_title = item.get("course_subtitle") or "Topic"
                s.shapes.title.text = _truncate_text(sub_title, 60)
                
                if len(s.placeholders) > 1:
                    body = s.placeholders[1]
                    content = item.get("subtitle_content") or ""
                    code_content = item.get("subtitle_code") or ""
                    help_text = item.get("subtitle_help_text") or ""
                    helpful_links = item.get("helpfull_links") or ""
                    
                    bullets = []
                    
                    # Add main content
                    if content:
                        content_lines = [ln.strip() for ln in content.split("\n") if ln.strip()]
                        for line in content_lines[:4]:
                            if len(line) > 10:  # Skip very short lines
                                bullets.append(f"📝 {_truncate_text(line, 120)}")
                    
                    # Add code examples if available
                    if code_content:
                        bullets.append("💻 Code Example:")
                        code_lines = [ln.strip() for ln in code_content.split("\n") if ln.strip()]
                        for line in code_lines[:3]:
                            bullets.append(f"  {_truncate_text(line, 100)}")
                    
                    # Add help text
                    if help_text:
                        bullets.append(f"💡 Tip: {_truncate_text(help_text, 100)}")
                    
                    # Add helpful links
                    if helpful_links:
                        bullets.append(f"🔗 Resources: {_truncate_text(helpful_links, 80)}")
                    
                    # Ensure we have at least some content
                    if not bullets:
                        bullets = [f"📚 {sub_title}", "Content will be added in future updates"]
                    
                    _add_bullets_to_placeholder(body, bullets[:8], font_size_pt=16)

        if slide_count < max_slides:
            # Last slide: Thanks
            s = prs.slides.add_slide(content_layout)
            s.shapes.title.text = "Thank You"
            if len(s.placeholders) > 1:
                body = s.placeholders[1]
                _add_bullets_to_placeholder(body, ["Questions?", f"Generated on {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"])

        # Generate filename if not provided
        if not filename:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"course_{course_id}_presentation_{timestamp}"
        
        # Create temp directory in project root (production-safe)
        project_root = os.path.dirname(os.path.dirname(os.path.dirname(__file__)))
        temp_dir = os.path.join(project_root, 'temp_files')
        os.makedirs(temp_dir, exist_ok=True)
        
        # Create temp file path
        temp_filename = f"{filename}_{datetime.now().strftime('%H%M%S')}.pptx"
        temp_path = os.path.join(temp_dir, temp_filename)
        
        # Save to temp file
        prs.save(temp_path)
        print(f"PPT saved locally to: {temp_path}")
        
        result = {
            'local_path': temp_path,
            'filename': f"{filename}.pptx",
            'cloud_url': None
        }
        
        # Upload to S3 if requested
        if upload_to_cloud:
            try:
                bucket_name = os.getenv('AWS_S3_BUCKET_NAME')
                if not bucket_name:
                    print("⚠️  AWS_S3_BUCKET_NAME not found in environment. Skipping S3 upload.")
                    return result
                
                # Create S3 object key
                s3_key = f"presentations/{filename}.pptx"
                
                print(f"Uploading to S3 bucket: {bucket_name}")
                cloud_url = upload_to_s3(temp_path, bucket_name, s3_key)
                result['cloud_url'] = cloud_url
                print(f"✅ PPT uploaded to S3: {cloud_url}")
                
                # Update database with PPT URL
                if update_course_ppt_url(conn, course_id, cloud_url):
                    print(f"📊 Course master table updated with PPT URL")
                else:
                    print(f"⚠️  Failed to update course master table")
                
                # Clean up local temp file after successful S3 upload
                try:
                    os.remove(temp_path)
                    print(f"🗑️  Local temp file cleaned up: {temp_filename}")
                    result['local_path'] = None  # File no longer exists locally
                except OSError as cleanup_error:
                    print(f"⚠️  Could not clean up temp file: {cleanup_error}")
                
            except Exception as e:
                print(f"❌ Failed to upload to S3: {str(e)}")
                print("PPT is still available locally")
                # Keep local file if S3 upload fails
        
        return result

    finally:
        if conn:
            conn.close()


